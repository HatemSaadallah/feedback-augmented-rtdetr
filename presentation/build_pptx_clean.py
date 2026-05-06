"""Clean, professional 18-slide TA presentation.

Run with:
    /home/hatem/miniconda3/envs/rtdetr/bin/python build_pptx_clean.py

Outputs:
    feedback_rtdetr_presentation.pptx     (the deliverable)
    _build/*.png                          (all rendered diagrams + math)

Design rules (per spec):
    - White background, navy / green / red palette only
    - Max 20–25 words per slide
    - Math slides: clean equation with labeled arrows pointing to variables
    - Big numbers for headline results (80–120pt)
    - Diagrams drawn in matplotlib so arrows land exactly where they should
"""

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle
from matplotlib import rcParams

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


ROOT = Path(__file__).parent
BUILD = ROOT / "_build"
BUILD.mkdir(exist_ok=True)
FIG_DIR = ROOT.parent / "report" / "figures"
V2_DIR = ROOT.parent / "v2_results"

OUT = ROOT / "feedback_rtdetr_presentation.pptx"

# ---- palette (navy + green + red on white) ----
NAVY  = RGBColor(0x1F, 0x6F, 0xEB)
DARK  = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT = RGBColor(0x66, 0x66, 0x77)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED   = RGBColor(0xD6, 0x2C, 0x2E)
BG    = RGBColor(0xFF, 0xFF, 0xFF)

NAVY_HEX  = "#1F6FEB"
DARK_HEX  = "#1A1A2E"
LIGHT_HEX = "#66667A"
GREEN_HEX = "#2CA02C"
RED_HEX   = "#D62C2E"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ===========================================================================
# Diagram + math rendering (matplotlib, transparent PNGs)
# ===========================================================================
def _save(fig, name):
    out = BUILD / name
    fig.savefig(out, dpi=300, bbox_inches="tight", pad_inches=0.05,
                transparent=True)
    plt.close(fig)
    return out


def fig_rtdetr_architecture():
    """RT-DETR architecture: backbone + hybrid encoder + deformable decoder."""
    fig, ax = plt.subplots(figsize=(13, 4.6))
    ax.axis("off"); ax.set_xlim(0, 13); ax.set_ylim(0, 4.6)

    # Image (left)
    img = FancyBboxPatch((0.3, 1.6), 1.5, 1.4, boxstyle="round,pad=0.05",
                          linewidth=2, edgecolor=DARK_HEX, facecolor="white")
    ax.add_patch(img)
    ax.text(1.05, 2.3, "Image", fontsize=15, ha="center", va="center",
            color=DARK_HEX, fontweight="bold")
    ax.text(1.05, 1.85, "640×640", fontsize=10, ha="center", color=LIGHT_HEX)

    # Backbone (ResNet-50-vd)
    bb = FancyBboxPatch((2.2, 0.9), 2.6, 2.8, boxstyle="round,pad=0.06",
                         linewidth=2.5, edgecolor=NAVY_HEX, facecolor="#F0F6FF")
    ax.add_patch(bb)
    ax.text(3.5, 3.3, "Backbone", fontsize=13, ha="center",
            color=NAVY_HEX, fontweight="bold")
    ax.text(3.5, 2.95, "ResNet-50-vd", fontsize=11, ha="center", color=DARK_HEX)
    # 4 pyramid output ticks
    for i, (level, stride) in enumerate([("P2", "4"), ("P3", "8"),
                                          ("P4", "16"), ("P5", "32")]):
        ax.text(3.5, 2.5 - i * 0.4,
                f"{level}  /  stride {stride}",
                fontsize=10, ha="center", color=LIGHT_HEX, family="monospace")

    # Hybrid Encoder (AIFI + CCFF)
    enc = FancyBboxPatch((5.2, 0.9), 2.7, 2.8, boxstyle="round,pad=0.06",
                          linewidth=2.5, edgecolor=NAVY_HEX, facecolor="#F0F6FF")
    ax.add_patch(enc)
    ax.text(6.55, 3.3, "Hybrid Encoder", fontsize=13, ha="center",
            color=NAVY_HEX, fontweight="bold")
    ax.text(6.55, 2.85, "AIFI", fontsize=11, ha="center", color=DARK_HEX,
            fontweight="bold")
    ax.text(6.55, 2.55, "(transformer on P5)", fontsize=9, ha="center",
            color=LIGHT_HEX, style="italic")
    ax.text(6.55, 2.05, "CCFF", fontsize=11, ha="center", color=DARK_HEX,
            fontweight="bold")
    ax.text(6.55, 1.75, "(cross-scale fusion)", fontsize=9, ha="center",
            color=LIGHT_HEX, style="italic")
    ax.text(6.55, 1.25, "→  memory  m", fontsize=10, ha="center",
            color=NAVY_HEX, fontweight="bold")

    # Decoder
    dec = FancyBboxPatch((8.3, 0.9), 2.7, 2.8, boxstyle="round,pad=0.06",
                          linewidth=2.5, edgecolor=NAVY_HEX, facecolor="#F0F6FF")
    ax.add_patch(dec)
    ax.text(9.65, 3.3, "Decoder", fontsize=13, ha="center",
            color=NAVY_HEX, fontweight="bold")
    ax.text(9.65, 2.85, "6 layers", fontsize=11, ha="center", color=DARK_HEX)
    ax.text(9.65, 2.45, "deformable", fontsize=11, ha="center", color=DARK_HEX)
    ax.text(9.65, 2.10, "cross-attention", fontsize=11, ha="center", color=DARK_HEX)
    ax.text(9.65, 1.55, "N object queries", fontsize=10, ha="center",
            color=LIGHT_HEX, style="italic")
    ax.text(9.65, 1.20, "Hungarian match → loss", fontsize=9, ha="center",
            color=LIGHT_HEX, style="italic")

    # Predictions
    out = FancyBboxPatch((11.3, 1.6), 1.5, 1.4, boxstyle="round,pad=0.05",
                         linewidth=2, edgecolor=DARK_HEX, facecolor="white")
    ax.add_patch(out)
    ax.text(12.05, 2.3, "Boxes +", fontsize=13, ha="center", va="center",
            color=DARK_HEX, fontweight="bold")
    ax.text(12.05, 1.95, "classes", fontsize=13, ha="center", va="center",
            color=DARK_HEX, fontweight="bold")

    # arrows
    for x_end, x_start in [(1.8, 2.2), (4.8, 5.2), (7.9, 8.3), (11.0, 11.3)]:
        ax.annotate("", xy=(x_start, 2.3), xytext=(x_end, 2.3),
                    arrowprops=dict(arrowstyle="->", lw=2.5, color=DARK_HEX))
    return _save(fig, "diag_rtdetr.png")


def fig_baseline_pipeline():
    """Image → Encoder → Decoder → Predictions, single-pass."""
    fig, ax = plt.subplots(figsize=(11, 2.4))
    ax.axis("off"); ax.set_xlim(0, 11); ax.set_ylim(0, 2.4)
    # 4 boxes — each box has its own (x, width); positions chosen so the
    # full diagram (including the rightmost box) fits inside [0, 11].
    boxes = [
        # (label,   x,    width)
        ("Image",   0.4, 1.6),
        ("Encoder", 2.7, 2.0),
        ("Decoder", 5.4, 2.0),
        ("Predictions", 8.1, 2.5),
    ]
    for label, x, w in boxes:
        rect = FancyBboxPatch((x, 0.7), w, 1.0, boxstyle="round,pad=0.05",
                              linewidth=2, edgecolor=NAVY_HEX, facecolor="white")
        ax.add_patch(rect)
        ax.text(x + w/2, 1.2, label, fontsize=20, ha="center", va="center",
                color=DARK_HEX, fontweight="bold")
    # arrows between consecutive boxes
    for (l1, x1, w1), (l2, x2, w2) in zip(boxes, boxes[1:]):
        ax.annotate("", xy=(x2, 1.2), xytext=(x1 + w1, 1.2),
                    arrowprops=dict(arrowstyle="->", lw=2.5, color=DARK_HEX))
    ax.text(5.5, 0.25, "one-way information flow",
            fontsize=14, ha="center", va="center",
            color=LIGHT_HEX, style="italic")
    return _save(fig, "diag_baseline.png")


def fig_idea_pipeline():
    """Encoder ↔ Decoder with feedback loop."""
    fig, ax = plt.subplots(figsize=(11, 3.2))
    ax.axis("off"); ax.set_xlim(0, 11); ax.set_ylim(0, 3.2)
    # 2 boxes
    enc = FancyBboxPatch((2.6, 1.2), 2.4, 1.0, boxstyle="round,pad=0.05",
                          linewidth=2.5, edgecolor=NAVY_HEX, facecolor="white")
    dec = FancyBboxPatch((6.4, 1.2), 2.4, 1.0, boxstyle="round,pad=0.05",
                          linewidth=2.5, edgecolor=NAVY_HEX, facecolor="white")
    ax.add_patch(enc); ax.add_patch(dec)
    ax.text(3.8, 1.7, "Encoder", fontsize=22, ha="center", va="center",
            color=DARK_HEX, fontweight="bold")
    ax.text(7.6, 1.7, "Decoder", fontsize=22, ha="center", va="center",
            color=DARK_HEX, fontweight="bold")
    # forward arrow (top)
    ax.annotate("", xy=(6.3, 1.85), xytext=(5.0, 1.85),
                arrowprops=dict(arrowstyle="->", lw=2.5, color=DARK_HEX))
    ax.text(5.65, 2.1, "memory", fontsize=12, ha="center", color=LIGHT_HEX)
    # feedback arrow (bottom, curved)
    arr = FancyArrowPatch((6.3, 1.55), (5.0, 1.55),
                          connectionstyle="arc3,rad=-0.5",
                          arrowstyle="->", lw=2.5, color=NAVY_HEX)
    ax.add_patch(arr)
    ax.text(5.65, 0.65, "feedback", fontsize=14, ha="center",
            color=NAVY_HEX, fontweight="bold")
    return _save(fig, "diag_idea.png")


def fig_attention_math():
    """Standard MHA equation centered, with three colored Q/K/V cards
    below acting as the legend. No arrows."""
    fig, ax = plt.subplots(figsize=(13, 5.6))
    ax.axis("off"); ax.set_xlim(0, 13); ax.set_ylim(0, 5.6)

    Q_COL = "#1F6FEB"
    K_COL = "#FF7F0E"
    V_COL = "#2CA02C"

    # equation, single render
    ax.text(6.5, 4.4,
            r"$\mathrm{Attn}(\mathbf{Q}, \mathbf{K}, \mathbf{V}) \;=\; "
            r"\mathrm{softmax}\!\left(\frac{\mathbf{Q}\,\mathbf{K}^{\top}}"
            r"{\sqrt{d}}\right)\mathbf{V}$",
            fontsize=42, ha="center", va="center", color=DARK_HEX)

    # 3 cards
    cards = [
        ("Q", "what we look for\n(query)",      Q_COL),
        ("K", "where we look\n(addresses)",     K_COL),
        ("V", "what we extract\n(contents)",    V_COL),
    ]
    for i, (letter, desc, col) in enumerate(cards):
        cx = 2.5 + i * 4.0
        # colored letter, large
        ax.text(cx, 2.5, letter, fontsize=48, ha="center", va="center",
                color=col, fontweight="bold")
        # description
        ax.text(cx, 1.55, desc, fontsize=14, ha="center", va="center",
                color=DARK_HEX)

    ax.text(6.5, 0.4,
            "Attention lets the model focus on important parts of the image.",
            fontsize=14, ha="center", va="center",
            color=LIGHT_HEX, style="italic")
    return _save(fig, "math_attention.png")


def fig_feedback_math():
    """Feedback formula centered, with three colored m/g/t cards below."""
    fig, ax = plt.subplots(figsize=(13, 5.6))
    ax.axis("off"); ax.set_xlim(0, 13); ax.set_ylim(0, 5.6)

    M_COL = "#1F6FEB"
    G_COL = "#FF7F0E"
    T_COL = "#2CA02C"

    # equation, single render
    ax.text(6.5, 4.4,
            r"$\mathbf{m}' \;=\; \mathbf{m} \;+\; g \cdot "
            r"\mathrm{Attn}(\mathbf{m},\, \mathbf{t})$",
            fontsize=48, ha="center", va="center", color=DARK_HEX)

    cards = [
        ("m", "encoder memory\n(image features)",     M_COL),
        ("g", "gate\n(how much feedback)",            G_COL),
        ("t", "decoder output\n(early predictions)",  T_COL),
    ]
    for i, (letter, desc, col) in enumerate(cards):
        cx = 2.5 + i * 4.0
        ax.text(cx, 2.5, letter, fontsize=48, ha="center", va="center",
                color=col, fontweight="bold")
        ax.text(cx, 1.55, desc, fontsize=14, ha="center", va="center",
                color=DARK_HEX)

    ax.text(6.5, 0.4,
            "We update the encoder using the decoder's own predictions.",
            fontsize=14, ha="center", va="center",
            color=LIGHT_HEX, style="italic")
    return _save(fig, "math_feedback.png")


def fig_v2_fixes():
    """Two boxes: gate floor + P2/P3 mask."""
    fig, ax = plt.subplots(figsize=(11, 3.6))
    ax.axis("off"); ax.set_xlim(0, 11); ax.set_ylim(0, 3.6)
    # box 1
    b1 = FancyBboxPatch((0.4, 0.6), 4.8, 2.4, boxstyle="round,pad=0.1",
                        linewidth=2.5, edgecolor=NAVY_HEX, facecolor="#F0F6FF")
    ax.add_patch(b1)
    ax.text(2.8, 2.55, "1.  Gate floor", fontsize=22, ha="center", va="center",
            color=NAVY_HEX, fontweight="bold")
    ax.text(2.8, 1.85, r"$g_\mathrm{eff} = 0.1 + 0.9\,\sigma(\alpha)$",
            fontsize=18, ha="center", va="center", color=DARK_HEX)
    ax.text(2.8, 1.05, "feedback can never be silenced",
            fontsize=14, ha="center", va="center", color=LIGHT_HEX, style="italic")
    # box 2
    b2 = FancyBboxPatch((5.8, 0.6), 4.8, 2.4, boxstyle="round,pad=0.1",
                        linewidth=2.5, edgecolor=GREEN_HEX, facecolor="#F0FFF0")
    ax.add_patch(b2)
    ax.text(8.2, 2.55, "2.  P2 / P3 only", fontsize=22, ha="center", va="center",
            color=GREEN_HEX, fontweight="bold")
    ax.text(8.2, 1.85, "refine where small objects live",
            fontsize=16, ha="center", va="center", color=DARK_HEX)
    ax.text(8.2, 1.05, "skip P4, P5 — focus the gradient",
            fontsize=14, ha="center", va="center", color=LIGHT_HEX, style="italic")
    return _save(fig, "diag_v2_fixes.png")


def fig_ablation_explainer():
    """Linear flow: checkpoint → 2 evals (stacked) → Δ result."""
    fig, ax = plt.subplots(figsize=(13, 4.2))
    ax.axis("off"); ax.set_xlim(0, 13); ax.set_ylim(0, 4.2)

    # left: the checkpoint
    cp = FancyBboxPatch((0.6, 1.6), 2.6, 1.0, boxstyle="round,pad=0.05",
                        linewidth=2.5, edgecolor=DARK_HEX, facecolor="white")
    ax.add_patch(cp)
    ax.text(1.9, 2.1, "v2 checkpoint", fontsize=18, ha="center", va="center",
            color=DARK_HEX, fontweight="bold")

    # arrows splitting up + down
    ax.annotate("", xy=(5.4, 3.3), xytext=(3.2, 2.4),
                arrowprops=dict(arrowstyle="->", lw=2.5, color=GREEN_HEX))
    ax.annotate("", xy=(5.4, 0.9), xytext=(3.2, 1.8),
                arrowprops=dict(arrowstyle="->", lw=2.5, color=RED_HEX))

    # two eval nodes
    on_box = FancyBboxPatch((5.4, 2.85), 3.4, 1.0, boxstyle="round,pad=0.05",
                             linewidth=2.5, edgecolor=GREEN_HEX,
                             facecolor="#E8F7E8")
    ax.add_patch(on_box)
    ax.text(7.1, 3.35, "evaluate · feedback ON", fontsize=15, ha="center",
            va="center", color=GREEN_HEX, fontweight="bold")

    off_box = FancyBboxPatch((5.4, 0.45), 3.4, 1.0, boxstyle="round,pad=0.05",
                              linewidth=2.5, edgecolor=RED_HEX,
                              facecolor="#FDECEC")
    ax.add_patch(off_box)
    ax.text(7.1, 0.95, "evaluate · feedback OFF", fontsize=15, ha="center",
            va="center", color=RED_HEX, fontweight="bold")

    # arrows merging into Δ
    ax.annotate("", xy=(11.2, 2.4), xytext=(8.8, 3.35),
                arrowprops=dict(arrowstyle="->", lw=2.5, color=GREEN_HEX))
    ax.annotate("", xy=(11.2, 1.8), xytext=(8.8, 0.95),
                arrowprops=dict(arrowstyle="->", lw=2.5, color=RED_HEX))

    # delta result on the right
    delta_box = FancyBboxPatch((11.0, 1.6), 1.6, 1.0, boxstyle="round,pad=0.05",
                                linewidth=2.5, edgecolor=NAVY_HEX,
                                facecolor="#F0F6FF")
    ax.add_patch(delta_box)
    ax.text(11.8, 2.1, "Δ", fontsize=32, ha="center", va="center",
            color=NAVY_HEX, fontweight="bold")
    return _save(fig, "diag_ablation.png")


def render_all_figures():
    rcParams["mathtext.fontset"] = "cm"
    print("Rendering matplotlib diagrams + math...")
    fig_rtdetr_architecture()
    fig_baseline_pipeline()
    fig_idea_pipeline()
    fig_attention_math()
    fig_feedback_math()
    fig_v2_fixes()
    fig_ablation_explainer()
    print(f"  → {len(list(BUILD.glob('*.png')))} PNGs in {BUILD}")


# ===========================================================================
# pptx helpers
# ===========================================================================
def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_baseline(run, percent):
    """Set sub/superscript via XML baseline shift (font-independent)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("baseline", str(int(percent)))


def add_runs(slide, x, y, w, h, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             default_size=18, default_color=DARK,
             default_font="Helvetica"):
    """Textbox with a list of styled runs.

    Each run is a dict with keys: text, size, bold, color, italic, font,
    subscript, superscript. Sub/superscript use real XML baseline shift
    (renders correctly even when the font lacks Unicode subscript glyphs).
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", default_font)
        run.font.size = Pt(r.get("size", default_size))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", default_color)
        if r.get("subscript"):
            _set_baseline(run, -25000)
        elif r.get("superscript"):
            _set_baseline(run, 30000)
    return tb


def aps_runs(size, color=DARK, bold=False, italic=False, prefix="", suffix=""):
    """Convenience: build runs for 'APₛ' with a real subscript S."""
    runs = []
    if prefix:
        runs.append({"text": prefix, "size": size, "color": color,
                     "bold": bold, "italic": italic})
    runs += [
        {"text": "AP", "size": size, "color": color,
         "bold": bold, "italic": italic},
        {"text": "S", "size": size, "color": color,
         "bold": bold, "italic": italic, "subscript": True},
    ]
    if suffix:
        runs.append({"text": suffix, "size": size, "color": color,
                     "bold": bold, "italic": italic})
    return runs


def apl_runs(size, color=DARK, bold=False, italic=False, prefix="", suffix=""):
    """Convenience: build runs for 'APₗ'."""
    runs = []
    if prefix:
        runs.append({"text": prefix, "size": size, "color": color,
                     "bold": bold, "italic": italic})
    runs += [
        {"text": "AP", "size": size, "color": color,
         "bold": bold, "italic": italic},
        {"text": "L", "size": size, "color": color,
         "bold": bold, "italic": italic, "subscript": True},
    ]
    if suffix:
        runs.append({"text": suffix, "size": size, "color": color,
                     "bold": bold, "italic": italic})
    return runs


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Helvetica", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_title(slide, text, *, color=NAVY):
    add_text(slide, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.9),
             text, size=34, bold=True, color=color)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(1.25), Inches(0.7), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def add_footer(slide, idx, total, section=None):
    left = "Feedback-Augmented RT-DETR  •  Hatem Saadallah & Nour Jennane"
    if section is not None:
        left = f"{section}  •  {left}"
    add_text(slide, Inches(0.7), Inches(7.05), Inches(10.0), Inches(0.3),
             left, size=10, color=LIGHT)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3),
             f"{idx} / {total}", size=10, color=LIGHT, align=PP_ALIGN.RIGHT)


def add_section_divider(prs, num, total_sections, title, subtitle):
    s = add_blank_slide(prs)
    # SECTION X / N (small caps top-left)
    add_text(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.4),
             f"SECTION  {num}  /  {total_sections}",
             size=14, bold=True, color=NAVY)
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(3.05), Inches(0.7), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    # section title (huge)
    add_text(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.2),
             title, size=54, bold=True, color=DARK)
    # subtitle (one-line preview of what's coming)
    add_text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.6),
             subtitle, size=20, color=LIGHT, italic=True)
    return s


# ===========================================================================
# Build the deck (rule-aligned 5-section structure)
# ===========================================================================
def build():
    on  = json.load(open(V2_DIR / "ablations" / "ablation_v2_feedback_on_640.json"))
    off = json.load(open(V2_DIR / "ablations" / "ablation_v2_feedback_off_640.json"))
    on800 = json.load(open(V2_DIR / "ablations" / "ablation_v2_feedback_on_800.json"))
    APS_ON  = on["AP_S"]  * 100
    APS_OFF = off["AP_S"] * 100
    APS_800 = on800["AP_S"] * 100
    DELTA = APS_ON - APS_OFF

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    TOTAL = 24
    SEC = "1. Problem Formulation"  # current section label, mutated below

    # ---------- 1. Title ----------
    s = add_blank_slide(prs)
    add_text(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.5),
             "BOCCONI UNIVERSITY  ·  COMPUTER VISION & IMAGE PROCESSING",
             size=12, color=NAVY, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(2.65), Inches(2.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    # Title — single line at 50pt (fits within 11.7" wide frame)
    add_text(s, Inches(0.8), Inches(2.85), Inches(11.7), Inches(1.1),
             "Feedback-Augmented RT-DETR",
             size=50, bold=True, color=DARK)
    # Subtitle on its own row, well below the title
    add_text(s, Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.7),
             "for Small Object Detection",
             size=28, color=LIGHT)
    # Authors + date
    add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
             "Hatem Saadallah  ·  Nour Jennane",
             size=18, color=DARK)
    add_text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.4),
             "Final Project  ·  April 2026",
             size=14, color=LIGHT)

    # ====================================================================
    # SECTION 1 — Problem Formulation
    # ====================================================================
    add_section_divider(prs, 1, 5, "Problem Formulation",
                        "What is small-object detection, and why does RT-DETR struggle?")
    add_footer(prs.slides[-1], 2, TOTAL, section="1. Problem Formulation")

    # ---------- 3. Problem ----------
    s = add_blank_slide(prs)
    add_title(s, "Small objects are harder to detect")
    # left: viz
    add_image(s, FIG_DIR / "viz_000000001000.png",
              Inches(0.7), Inches(1.7), w=Inches(7.0))
    # right: numbers
    add_runs(s, Inches(8.4), Inches(2.0), Inches(4.5), Inches(0.5),
             apl_runs(15, color=LIGHT, suffix=" (large)"))
    add_text(s, Inches(8.4), Inches(2.4), Inches(4.5), Inches(1.4),
             "67.7", size=84, bold=True, color=GREEN)
    add_runs(s, Inches(8.4), Inches(4.0), Inches(4.5), Inches(0.5),
             aps_runs(15, color=LIGHT, suffix=" (small)"))
    add_text(s, Inches(8.4), Inches(4.4), Inches(4.5), Inches(1.4),
             "34.7", size=84, bold=True, color=RED)
    add_text(s, Inches(8.4), Inches(6.0), Inches(4.5), Inches(0.6),
             "33-point gap. Why?",
             size=18, color=DARK, italic=True)
    add_footer(s, 3, TOTAL, section="1. Problem Formulation")

    # ---------- 4. RT-DETR architecture ----------
    s = add_blank_slide(prs)
    add_title(s, "RT-DETR: backbone → encoder → decoder")
    add_image(s, BUILD / "diag_rtdetr.png",
              Inches(0.4), Inches(1.6), w=Inches(12.5))
    add_text(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(0.5),
             "End-to-end, NMS-free. Six decoder layers all attend over the same encoder memory.",
             size=18, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.4),
             "Backbone extracts a feature pyramid (P2–P5); the hybrid encoder fuses scales; the decoder cross-attends.",
             size=12, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 4, TOTAL, section="1. Problem Formulation")

    # ---------- 5. Why? structural one-way ----------
    s = add_blank_slide(prs)
    add_title(s, "The reason is structural")
    add_image(s, BUILD / "diag_baseline.png",
              Inches(0.8), Inches(2.5), w=Inches(11.7))
    add_text(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.6),
             "Encoder memory is computed once. The decoder cannot revise it.",
             size=22, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.5),
             "Refinement of the features small objects live in is a one-way street.",
             size=14, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 5, TOTAL, section="1. Problem Formulation")

    # ====================================================================
    # SECTION 2 — Data Sourcing Strategy
    # ====================================================================
    add_section_divider(prs, 2, 5, "Data Sourcing Strategy",
                        "What dataset we used and why it is appropriate.")
    add_footer(prs.slides[-1], 6, TOTAL, section="2. Data Sourcing Strategy")

    # ---------- 6. COCO 2017 ----------
    s = add_blank_slide(prs)
    add_title(s, "COCO 2017 — the standard small-object benchmark")
    # left card: dataset facts
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.7), Inches(1.7), Inches(6.3), Inches(4.6))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)
    card.line.color.rgb = LIGHT; card.line.width = Pt(0.75)
    card.shadow.inherit = False
    add_text(s, Inches(0.95), Inches(1.95), Inches(5.85), Inches(0.5),
             "Dataset", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.95), Inches(2.45), Inches(5.85), Inches(0.5),
             "Microsoft COCO 2017", size=24, bold=True, color=DARK)
    add_text(s, Inches(0.95), Inches(3.1), Inches(5.85), Inches(0.5),
             "Splits", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.95), Inches(3.55), Inches(5.85), Inches(0.5),
             "train2017:  118k images, 860k boxes",
             size=15, color=DARK, font="Menlo")
    add_text(s, Inches(0.95), Inches(3.95), Inches(5.85), Inches(0.5),
             "val2017:        5k images   (held out)",
             size=15, color=DARK, font="Menlo")
    add_text(s, Inches(0.95), Inches(4.7), Inches(5.85), Inches(0.5),
             "Pre-trained weights", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.95), Inches(5.15), Inches(5.85), Inches(0.5),
             "rtdetr_r50vd_6x_coco.pth  (public, finetuned 13 ep)",
             size=13, color=DARK, font="Menlo")
    # right card: why
    card2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(7.4), Inches(1.7), Inches(5.3), Inches(4.6))
    card2.fill.solid(); card2.fill.fore_color.rgb = RGBColor(0xE8, 0xF1, 0xFF)
    card2.line.color.rgb = NAVY; card2.line.width = Pt(1.5)
    card2.shadow.inherit = False
    add_text(s, Inches(7.6), Inches(1.95), Inches(5.0), Inches(0.5),
             "Why COCO?", size=18, bold=True, color=NAVY)
    why = [
        ("Standard benchmark — directly comparable to RT-DETR, DETR, YOLO.", None),
        ("Has explicit small / medium / large size buckets (S / M / L).", None),
        ("~41% of annotated boxes are small (area < 32² px).", None),
        ("Same pre-training corpus as our baseline → fair comparison.", None),
    ]
    for i, (line, _) in enumerate(why):
        add_text(s, Inches(7.6), Inches(2.6 + i * 0.85), Inches(5.0), Inches(0.8),
                 "•  " + line, size=14, color=DARK)
    add_footer(s, 7, TOTAL, section="2. Data Sourcing Strategy")

    # ====================================================================
    # SECTION 3 — Proposed Solution
    # ====================================================================
    add_section_divider(prs, 3, 5, "Proposed Solution",
                        "The feedback module, the math behind it, and the two fixes that make it work.")
    add_footer(prs.slides[-1], 8, TOTAL, section="3. Proposed Solution")

    # ---------- 8. Our idea ----------
    s = add_blank_slide(prs)
    add_title(s, "Our idea: let the model refine its own features")
    add_image(s, BUILD / "diag_idea.png",
              Inches(0.8), Inches(2.2), w=Inches(11.7))
    add_text(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.6),
             "Feed early decoder predictions back into the encoder memory.",
             size=22, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 9, TOTAL, section="3. Proposed Solution")

    # ---------- 10. MATH 1: Attention ----------
    s = add_blank_slide(prs)
    add_title(s, "Attention, in one line")
    add_image(s, BUILD / "math_attention.png",
              Inches(0.5), Inches(1.5), w=Inches(12.3))
    add_footer(s, 10, TOTAL, section="3. Proposed Solution")

    # ---------- 11. MATH 2: Feedback formula ----------
    s = add_blank_slide(prs)
    add_title(s, "Our feedback rule")
    add_image(s, BUILD / "math_feedback.png",
              Inches(0.5), Inches(1.5), w=Inches(12.3))
    add_footer(s, 11, TOTAL, section="3. Proposed Solution")

    # ---------- 12. v1 first attempt + why it failed (combined) ----------
    s = add_blank_slide(prs)
    add_title(s, "First attempt (v1) — and why it failed")
    add_text(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.5),
             "Plain sigmoid gate, applied to all four pyramid levels.",
             size=18, color=DARK, italic=True, align=PP_ALIGN.CENTER)
    # left: the result (real subscript via runs)
    add_runs(s, Inches(0.4), Inches(2.8), Inches(6.4), Inches(1.6),
             [{"text": "Δ AP", "size": 56, "bold": True, "color": RED},
              {"text": "S",    "size": 56, "bold": True, "color": RED, "subscript": True},
              {"text": "  =  0.00", "size": 56, "bold": True, "color": RED}],
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.4), Inches(4.5), Inches(6.4), Inches(0.5),
             "the mechanism contributed nothing",
             size=15, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
    # right: explanation
    add_text(s, Inches(7.1), Inches(2.5), Inches(5.6), Inches(0.6),
             "Why?", size=24, bold=True, color=DARK)
    add_text(s, Inches(7.1), Inches(3.1), Inches(5.6), Inches(2.5),
             "α drifts to −∞ → gate ≈ 0.\n\nThe feedback signal is multiplied by zero before reaching memory.",
             size=15, color=DARK)
    add_footer(s, 12, TOTAL, section="3. Proposed Solution")

    # ---------- 13. v2 fixes ----------
    s = add_blank_slide(prs)
    add_title(s, "v2: two fixes — zero new parameters")
    add_image(s, BUILD / "diag_v2_fixes.png",
              Inches(1.0), Inches(2.2), w=Inches(11.3))
    add_text(s, Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.5),
             "Reparameterize the gate. Refine only the levels small objects live in.",
             size=18, color=DARK, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 13, TOTAL, section="3. Proposed Solution")

    # ====================================================================
    # SECTION 4 — Performance Evaluation Approach
    # ====================================================================
    add_section_divider(prs, 4, 5, "Performance Evaluation Approach",
                        "Metrics, baselines, and the same-checkpoint ablation that isolates causality.")
    add_footer(prs.slides[-1], 14, TOTAL, section="4. Evaluation Approach")

    # ---------- 15. Metrics + protocol ----------
    s = add_blank_slide(prs)
    add_title(s, "How we evaluate the contribution")
    # top-left: metrics card (compacted)
    card_m = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.7), Inches(1.55), Inches(6.0), Inches(1.85))
    card_m.fill.solid(); card_m.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)
    card_m.line.color.rgb = LIGHT; card_m.line.width = Pt(0.75)
    card_m.shadow.inherit = False
    add_text(s, Inches(0.95), Inches(1.7), Inches(5.6), Inches(0.4),
             "Metrics  (COCO standard)", size=13, bold=True, color=NAVY)
    add_runs(s, Inches(0.95), Inches(2.1), Inches(5.6), Inches(0.4),
             [{"text": "AP, AP", "size": 13, "color": DARK},
              {"text": "S",      "size": 13, "color": DARK, "subscript": True},
              {"text": ", AP",   "size": 13, "color": DARK},
              {"text": "M",      "size": 13, "color": DARK, "subscript": True},
              {"text": ", AP",   "size": 13, "color": DARK},
              {"text": "L",      "size": 13, "color": DARK, "subscript": True},
              {"text": "  —  IoU 0.5..0.95", "size": 13, "color": DARK}])
    add_runs(s, Inches(0.95), Inches(2.5), Inches(5.6), Inches(0.4),
             [{"text": "AP", "size": 13, "color": DARK, "bold": True},
              {"text": "S",  "size": 13, "color": DARK, "bold": True, "subscript": True},
              {"text": " is the headline metric (area < 32² px)",
               "size": 13, "color": DARK, "bold": True}])
    add_text(s, Inches(0.95), Inches(2.9), Inches(5.6), Inches(0.4),
             "Others verify we don't degrade easier categories.",
             size=10, color=LIGHT, italic=True)
    # top-right: baselines card (compacted)
    card_b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(7.0), Inches(1.55), Inches(5.6), Inches(1.85))
    card_b.fill.solid(); card_b.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)
    card_b.line.color.rgb = LIGHT; card_b.line.width = Pt(0.75)
    card_b.shadow.inherit = False
    add_text(s, Inches(7.25), Inches(1.7), Inches(5.2), Inches(0.4),
             "Baselines", size=13, bold=True, color=NAVY)
    add_text(s, Inches(7.25), Inches(2.1), Inches(5.2), Inches(0.4),
             "RT-DETR-R50  (published)",
             size=13, color=DARK)
    add_text(s, Inches(7.25), Inches(2.5), Inches(5.2), Inches(0.4),
             "v1 feedback  (before the fix)",
             size=13, color=DARK)
    add_text(s, Inches(7.25), Inches(2.9), Inches(5.2), Inches(0.4),
             "v2 feedback  (with the fix)",
             size=13, color=DARK, bold=True)
    # bottom: the ablation diagram (centered, smaller)
    add_text(s, Inches(0.7), Inches(3.65), Inches(12.0), Inches(0.5),
             "Same-checkpoint ablation  —  isolates causal contribution",
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_image(s, BUILD / "diag_ablation.png",
              Inches(2.4), Inches(4.2), w=Inches(8.5))
    add_text(s, Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.4),
             "Toggle feedback ON / OFF at inference. Identical weights, identical inputs.",
             size=12, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 15, TOTAL, section="4. Evaluation Approach")

    # ====================================================================
    # SECTION 5 — Final Results
    # ====================================================================
    add_section_divider(prs, 5, 5, "Final Results",
                        "Training trajectory, the causal effect, robustness, trade-offs, and what's next.")
    add_footer(prs.slides[-1], 16, TOTAL, section="5. Final Results")

    # ---------- 17. Training trajectory ----------
    s = add_blank_slide(prs)
    add_title(s, "Training: v2 climbs faster, ends higher")
    add_image(s, FIG_DIR / "learning_curves.png",
              Inches(2.4), Inches(1.6), w=Inches(8.5))
    add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.4),
             "v2 first crosses v1's final 33.91 at epoch 8 — four epochs early.",
             size=14, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 17, TOTAL, section="5. Final Results")

    # ---------- 18. Final scores (v1 vs v2 with feedback ON) ----------
    s = add_blank_slide(prs)
    add_title(s, "Final performance (with feedback ON)")
    # left card v1
    c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.0), Inches(2.2), Inches(5.3), Inches(3.5))
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)
    c1.line.color.rgb = LIGHT; c1.line.width = Pt(0.75)
    c1.shadow.inherit = False
    add_text(s, Inches(1.0), Inches(2.4), Inches(5.3), Inches(0.5),
             "v1", size=24, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(2.95), Inches(5.3), Inches(2.0),
             "33.91", size=88, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    add_runs(s, Inches(1.0), Inches(5.0), Inches(5.3), Inches(0.5),
             aps_runs(18, color=LIGHT, italic=True),
             align=PP_ALIGN.CENTER)
    # right card v2
    c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(7.0), Inches(2.2), Inches(5.3), Inches(3.5))
    c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0xE8, 0xF7, 0xE8)
    c2.line.color.rgb = GREEN; c2.line.width = Pt(2)
    c2.shadow.inherit = False
    add_text(s, Inches(7.0), Inches(2.4), Inches(5.3), Inches(0.5),
             "v2", size=24, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(7.0), Inches(2.95), Inches(5.3), Inches(2.0),
             "34.30", size=88, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_runs(s, Inches(7.0), Inches(5.0), Inches(5.3), Inches(0.5),
             aps_runs(18, color=GREEN, italic=True, suffix="    +0.4"),
             align=PP_ALIGN.CENTER)
    add_footer(s, 18, TOTAL, section="5. Final Results")

    # ---------- 19. Causal result (the big one) ----------
    s = add_blank_slide(prs)
    add_title(s, "Causal effect of feedback")
    cols = [
        ("ON",  f"{APS_ON:.2f}",  GREEN, "feedback enabled"),
        ("OFF", f"{APS_OFF:.2f}", LIGHT, "feedback bypassed"),
        ("Δ",   f"+{DELTA:.2f}",  NAVY,  "the contribution"),
    ]
    for i, (label, val, c, sub) in enumerate(cols):
        x = Inches(0.7 + i * 4.15)
        add_text(s, x, Inches(2.0), Inches(4.0), Inches(0.6),
                 label, size=24, color=LIGHT, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.7), Inches(4.0), Inches(2.2),
                 val, size=110, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(5.2), Inches(4.0), Inches(0.5),
                 sub, size=14, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.6),
             "Feedback has a real causal impact on small-object detection.",
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, 19, TOTAL, section="5. Final Results")

    # ---------- 20. Resolution effect ----------
    s = add_blank_slide(prs)
    add_title(s, "Higher resolution → bigger gain")
    add_text(s, Inches(0.7), Inches(2.4), Inches(5.8), Inches(0.5),
             "640 × 640", size=22, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(3.0), Inches(5.8), Inches(2.0),
             f"{APS_ON:.2f}", size=88, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                              Inches(5.8), Inches(3.6), Inches(1.7), Inches(0.6))
    arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
    arr.line.fill.background()
    add_text(s, Inches(6.8), Inches(2.4), Inches(5.8), Inches(0.5),
             "800 × 800", size=22, color=NAVY, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(6.8), Inches(3.0), Inches(5.8), Inches(2.0),
             f"{APS_800:.2f}", size=88, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(0.5),
             f"+{APS_800 - APS_ON:.2f} from more pixels.",
             size=22, color=DARK, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.5),
             "Feedback helps. Spatial sampling helps too.",
             size=16, color=LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, 20, TOTAL, section="5. Final Results")

    # ---------- 21. Strengths + Limitations (combined) ----------
    s = add_blank_slide(prs)
    add_title(s, "Strengths and limitations")
    # green +
    plus = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.7), Inches(1.8), Inches(5.9), Inches(4.7))
    plus.fill.solid(); plus.fill.fore_color.rgb = RGBColor(0xE8, 0xF7, 0xE8)
    plus.line.color.rgb = GREEN; plus.line.width = Pt(1.5)
    plus.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(2.0), Inches(5.5), Inches(0.6),
             "+  Strengths", size=22, bold=True, color=GREEN)
    pluses = [
        "Causal contribution: +0.99 on small-object AP, same checkpoint.",
        "Zero new parameters relative to v1.",
        "Consistent improvement on every metric (AP, S, M, L).",
        "Robust at higher resolution (37.01 small-object AP at 800×800).",
    ]
    for i, line in enumerate(pluses):
        add_text(s, Inches(0.95), Inches(2.7 + i * 0.85), Inches(5.5), Inches(0.8),
                 "•  " + line, size=14, color=DARK)
    # red −
    minus = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.7))
    minus.fill.solid(); minus.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xEC)
    minus.line.color.rgb = RED; minus.line.width = Pt(1.5)
    minus.shadow.inherit = False
    add_text(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.6),
             "−  Limitations", size=22, bold=True, color=RED)
    minuses = [
        "Single seed — multi-seed variance unmeasured.",
        "13-epoch finetune (vs 72-epoch from scratch).",
        "Floor and mask not isolated (a 2×2 ablation is missing).",
        "≈2× slower than baseline (53 → 26 FPS) — cost is the P2 level.",
    ]
    for i, line in enumerate(minuses):
        add_text(s, Inches(7.05), Inches(2.7 + i * 0.85), Inches(5.5), Inches(0.8),
                 "•  " + line, size=14, color=DARK)
    add_text(s, Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.4),
             "Strong on accuracy; pays in inference speed. The next slide proposes a way to recover speed.",
             size=12, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 21, TOTAL, section="5. Final Results")

    # ---------- 22. Future work — recovering FPS (NEW) ----------
    s = add_blank_slide(prs)
    add_title(s, "Future work: recovering inference speed")
    add_text(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.5),
             "The 2× slowdown comes from the P2 level, not the feedback module.",
             size=18, color=DARK, italic=True, align=PP_ALIGN.CENTER)
    # 3 directions, each with a How: line and a Goal: line
    directions = [
        ("Train with P2, infer without",
         "fastest path",
         "Use P2 during training to improve small-object learning, then remove it at inference.",
         "Keep most of the +0.99 small-object AP gain without the P2 computational cost."),
        ("Lighter P2 path",
         "compromise",
         "Reduce the cost of P2 — fewer channels, or a simpler projection layer.",
         "Preserve small-object improvements while reducing the high-resolution P2 slowdown."),
        ("Knowledge distillation",
         "principled",
         "Transfer v2's small-object improvements into a faster baseline RT-DETR student.",
         "Recover baseline inference speed while retaining the small-object AP gains learned from P2 feedback."),
    ]
    for i, (head, tag, how, goal) in enumerate(directions):
        x = Inches(0.7 + i * 4.15)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(2.4), Inches(4.0), Inches(4.0))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF0, 0xF6, 0xFF)
        card.line.color.rgb = NAVY; card.line.width = Pt(1.2)
        card.shadow.inherit = False
        # number badge
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   x + Inches(0.2), Inches(2.55),
                                   Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = NAVY
        circ.line.fill.background()
        add_text(s, x + Inches(0.2), Inches(2.58), Inches(0.5), Inches(0.5),
                 str(i + 1), size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
        # title
        add_text(s, x + Inches(0.85), Inches(2.55), Inches(3.0), Inches(0.6),
                 head, size=14, bold=True, color=DARK)
        # tag
        add_text(s, x + Inches(0.2), Inches(3.25), Inches(3.6), Inches(0.4),
                 tag, size=10, color=NAVY, italic=True, bold=True)
        # how
        add_text(s, x + Inches(0.2), Inches(3.75), Inches(3.6), Inches(0.4),
                 "How", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), Inches(4.10), Inches(3.6), Inches(1.4),
                 how, size=11, color=DARK)
        # goal
        add_text(s, x + Inches(0.2), Inches(5.30), Inches(3.6), Inches(0.4),
                 "Goal", size=11, bold=True, color=GREEN)
        add_text(s, x + Inches(0.2), Inches(5.65), Inches(3.6), Inches(1.4),
                 goal, size=11, color=DARK)
    add_text(s, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
             "Common goal: keep the +0.99 small-object AP gain while approaching the 53 FPS baseline.",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, 22, TOTAL, section="5. Final Results")

    # ---------- 23. Conclusion ----------
    s = add_blank_slide(prs)
    add_title(s, "Conclusion")
    bullets = [
        ("Feedback improves small-object detection",
         "+0.99 small-object AP, causal contribution at inference."),
        ("Only works with proper design",
         "Floor the gate; restrict to small-object levels."),
        ("Modest but real",
         "Consistent across metrics; reproducible from the same checkpoint."),
    ]
    for i, (head, sub) in enumerate(bullets):
        y = Inches(2.0 + i * 1.55)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(0.9), y + Inches(0.15),
                                   Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = NAVY
        circ.line.fill.background()
        add_text(s, Inches(0.9), y + Inches(0.18), Inches(0.55), Inches(0.55),
                 str(i + 1), size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.7), y + Inches(0.05), Inches(11.0), Inches(0.6),
                 head, size=24, bold=True, color=DARK)
        add_text(s, Inches(1.7), y + Inches(0.7), Inches(11.0), Inches(0.6),
                 sub, size=15, color=LIGHT, italic=True)
    add_footer(s, 23, TOTAL, section="5. Final Results")

    # ---------- 24. Thank you ----------
    s = add_blank_slide(prs)
    add_text(s, Inches(0.7), Inches(2.4), Inches(12.0), Inches(1.5),
             "Thank you.",
             size=80, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.2), Inches(12.0), Inches(0.8),
             "Questions?",
             size=40, color=NAVY, align=PP_ALIGN.CENTER)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(6.17), Inches(5.4), Inches(1.0), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.4),
             "Hatem Saadallah  ·  Nour Jennane",
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.4),
             "github.com/HatemSaadallah/feedback-augmented-rtdetr",
             size=12, color=LIGHT, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"\nSaved: {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, {TOTAL} slides)")


if __name__ == "__main__":
    render_all_figures()
    build()

"""Build the TA presentation deck.

Run with:
    /home/hatem/miniconda3/envs/rtdetr/bin/python build_pptx.py

Outputs:
    feedback_augmented_rtdetr.pptx
    equations/*.png  (high-DPI math, regenerated each run)

Slides are intentionally sparse: big visuals, math when math is the point,
and hyperlinks to the GitHub repo + report PDF as the interactive element.
"""

import json
import os
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

ROOT = Path(__file__).parent
EQ_DIR = ROOT / "equations"
FIG_DIR = ROOT.parent / "report" / "figures"
V1_DIR = ROOT.parent / "v1_results"
V2_DIR = ROOT.parent / "v2_results"
OUT = ROOT / "feedback_augmented_rtdetr.pptx"

# ----- color palette -----
NAVY = RGBColor(0x1F, 0x6F, 0xEB)   # title / accent
DARK = RGBColor(0x1A, 0x1A, 0x2E)   # body text
LIGHT = RGBColor(0x55, 0x55, 0x66)  # secondary text
GREEN = RGBColor(0x2C, 0xA0, 0x2C)  # positive result
RED = RGBColor(0xD6, 0x2C, 0x2E)    # negative result
BG = RGBColor(0xFF, 0xFF, 0xFF)     # slide background

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ===========================================================================
# 1) Generate high-quality math equation PNGs via matplotlib mathtext
# ===========================================================================
def render_eq(latex, filename, fontsize=42, color="#1A1A2E"):
    """Render a LaTeX math expression as a transparent high-DPI PNG."""
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.patch.set_alpha(0.0)
    fig.text(0.5, 0.5, latex, fontsize=fontsize, color=color,
             ha="center", va="center")
    out = EQ_DIR / filename
    fig.savefig(out, dpi=300, bbox_inches="tight", pad_inches=0.05,
                transparent=True)
    plt.close(fig)
    return out


def generate_equations():
    EQ_DIR.mkdir(exist_ok=True)
    # Plain mathtext (no usetex) — supports a broad subset of LaTeX without
    # invoking the external pdflatex binary. We avoid \underbrace etc.
    rcParams["mathtext.fontset"] = "cm"
    rcParams["mathtext.rm"] = "serif"

    eqs = {
        # The headline gate equation (v2 contribution)
        "gate_floor.png":
            r"$g_\mathrm{eff} \;=\; \mathrm{floor} \;+\; "
            r"(1 - \mathrm{floor}) \cdot \sigma(\alpha)$",
        # Standard MHA attention (background)
        "attn.png":
            r"$\mathrm{Attn}(Q, K, V) \;=\; \mathrm{softmax}\!\left("
            r"\frac{Q K^\top}{\sqrt{d_k}}\right) V$",
        # The feedback module
        "feedback_module.png":
            r"$m' = \mathrm{LN}\!\left(m \;+\; g_\mathrm{eff} \cdot "
            r"\mathrm{MHA}(m,\, t^{(1)},\, t^{(1)})\right)$",
        # Memory shape — without \underbrace (matplotlib mathtext doesn't support it)
        "memory_shape.png":
            r"$L \;=\; \underset{\mathrm{P2}}{160{\times}160} "
            r"\;+\; \underset{\mathrm{P3}}{80{\times}80} "
            r"\;+\; \underset{\mathrm{P4}}{40{\times}40} "
            r"\;+\; \underset{\mathrm{P5}}{20{\times}20} "
            r"\;=\; 34{,}000$",
        # The headline result
        "delta.png":
            r"$\Delta\mathrm{AP}_S = \mathrm{AP}_S^\mathrm{ON} - "
            r"\mathrm{AP}_S^\mathrm{OFF} = +0.99$",
        # v1 silent result — prefix the version label rather than superscript
        "v1_silent.png":
            r"$\mathrm{v1}:\quad \Delta\mathrm{AP}_S \;=\; 0.00$",
        # v2 result emphasis
        "v2_pass.png":
            r"$\mathrm{v2}:\quad \Delta\mathrm{AP}_S \;=\; +0.99$",
        # Level mask
        "level_mask.png":
            r"$\mathbf{b} = (1, 1, 0, 0) \;\Rightarrow\; "
            r"\mathrm{refine}\;\{\mathrm{P2}, \mathrm{P3}\}$",
    }

    print("Generating math equations...")
    for filename, latex in eqs.items():
        # Render positive/negative variants in proper colors
        if filename == "v1_silent.png":
            render_eq(latex, filename, color="#D62C2E", fontsize=56)
        elif filename in ("v2_pass.png", "delta.png"):
            render_eq(latex, filename, color="#2CA02C", fontsize=56)
        else:
            render_eq(latex, filename, fontsize=44)
        print(f"  {filename}")


# ===========================================================================
# 2) Slide-builder helpers
# ===========================================================================
def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Helvetica"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_title(slide, text):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7),
             text, size=32, bold=True, color=NAVY)
    # underline accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.05), Inches(0.6), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def add_hyperlink(slide, x, y, w, h, label, url, *, size=14):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.name = "Helvetica"
    run.font.size = Pt(size)
    run.font.bold = False
    run.font.color.rgb = NAVY
    run.font.underline = True
    run.hyperlink.address = url
    return tb


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8.0), Inches(0.3),
             "Feedback-Augmented RT-DETR  •  Hatem Saadallah & Nour Jennane  •  Bocconi University",
             size=10, color=LIGHT)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3),
             f"{idx} / {total}", size=10, color=LIGHT, align=PP_ALIGN.RIGHT)


def load_v2_numbers():
    """Pull the headline numbers fresh from the JSONs."""
    on = json.load(open(V2_DIR / "ablations" / "ablation_v2_feedback_on_640.json"))
    off = json.load(open(V2_DIR / "ablations" / "ablation_v2_feedback_off_640.json"))
    on800 = json.load(open(V2_DIR / "ablations" / "ablation_v2_feedback_on_800.json"))
    return {
        "AP_on": on["AP"] * 100, "AP_off": off["AP"] * 100,
        "APS_on": on["AP_S"] * 100, "APS_off": off["AP_S"] * 100,
        "APM_on": on["AP_M"] * 100, "APM_off": off["AP_M"] * 100,
        "APL_on": on["AP_L"] * 100, "APL_off": off["AP_L"] * 100,
        "APS_800": on800["AP_S"] * 100,
    }


# ===========================================================================
# 3) Build the deck
# ===========================================================================
def build():
    nums = load_v2_numbers()
    delta_aps = nums["APS_on"] - nums["APS_off"]
    print(f"Loaded ΔAP_S = {delta_aps:+.2f}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 10

    # ---------- Slide 1: Title ----------
    s = add_blank_slide(prs)
    add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.6),
             "COMPUTER VISION AND IMAGE PROCESSING — BOCCONI UNIVERSITY",
             size=12, color=NAVY, bold=True)
    add_text(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.4),
             "Feedback-Augmented RT-DETR",
             size=54, bold=True, color=DARK)
    add_text(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.7),
             "A Cross-Attention Refinement Strategy for Enhanced Small-Object Detection",
             size=22, color=LIGHT)
    add_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
             "Hatem Saadallah  ·  Nour Jennane",
             size=18, color=DARK)
    add_text(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.4),
             "Final Project · April 2026",
             size=14, color=LIGHT)
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(2.95), Inches(2.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # ---------- Slide 2: The problem ----------
    s = add_blank_slide(prs)
    add_title(s, "The problem: small objects in RT-DETR")
    add_text(s, Inches(0.6), Inches(1.6), Inches(7.5), Inches(0.6),
             "RT-DETR is fast and NMS-free — but small-object AP lags.",
             size=22, color=DARK)
    # left side: three big numbers
    add_text(s, Inches(0.6), Inches(2.6), Inches(3.0), Inches(0.4),
             "AP_S (small)", size=14, color=LIGHT)
    add_text(s, Inches(0.6), Inches(2.95), Inches(3.0), Inches(0.9),
             "34.7", size=64, bold=True, color=RED)
    add_text(s, Inches(4.0), Inches(2.6), Inches(3.0), Inches(0.4),
             "AP_L (large)", size=14, color=LIGHT)
    add_text(s, Inches(4.0), Inches(2.95), Inches(3.0), Inches(0.9),
             "67.7", size=64, bold=True, color=GREEN)
    add_text(s, Inches(0.6), Inches(4.0), Inches(7.4), Inches(0.5),
             "33-point gap on the same model. Why?",
             size=18, color=DARK)
    add_text(s, Inches(0.6), Inches(4.7), Inches(7.4), Inches(2.3),
             ("Encoder memory is computed once before any decoder layer fires.\n"
              "It is never refined as the decoder accumulates information.\n"
              "Small objects need every available cue at every layer."),
             size=16, color=LIGHT)
    # right side: detection viz
    add_image(s, FIG_DIR / "viz_000000001000.png",
              Inches(8.3), Inches(1.6), w=Inches(4.6))
    add_footer(s, 2, TOTAL)

    # ---------- Slide 3: The idea ----------
    s = add_blank_slide(prs)
    add_title(s, "The idea: feed early predictions back into memory")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.6),
             "After decoder layer 1, refine encoder memory using its own predictions.",
             size=20, color=DARK)
    add_image(s, EQ_DIR / "feedback_module.png",
              Inches(2.0), Inches(2.6), w=Inches(9.3))
    add_text(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(0.5),
             "Cross-attention with a learnable scalar gate g_eff.",
             size=16, color=LIGHT, align=PP_ALIGN.CENTER)
    add_image(s, EQ_DIR / "memory_shape.png",
              Inches(0.8), Inches(5.4), w=Inches(11.7))
    add_text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
             "75% of memory tokens live at stride 4 — the natural resolution for small objects.",
             size=13, color=LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, 3, TOTAL)

    # ---------- Slide 4: v1 — the negative result ----------
    s = add_blank_slide(prs)
    add_title(s, "v1: the mechanism trained — but contributed nothing")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.6),
             "Same-checkpoint inference ablation on COCO val2017 @ 640:",
             size=20, color=DARK)
    add_image(s, EQ_DIR / "v1_silent.png",
              Inches(2.0), Inches(2.8), w=Inches(9.3))
    add_text(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(0.6),
             "Toggling feedback on/off at inference made zero difference.",
             size=18, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.6),
             "The cross-attention weights were trained, but the gate had decayed to ≈0.12 during training.\n"
             "The feedback signal was being multiplied by ~zero before reaching memory.",
             size=14, color=LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, 4, TOTAL)

    # ---------- Slide 5: Diagnosis — gate collapse ----------
    s = add_blank_slide(prs)
    add_title(s, "Diagnosis: the gate had an escape route")
    add_text(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.6),
             "Plain sigmoid gate.",
             size=22, bold=True, color=DARK)
    add_text(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(2.5),
             ("• Gradient → 0 as α → −∞\n\n"
              "• If the rest of the network can compensate, the optimizer\n  pushes α arbitrarily negative at no cost.\n\n"
              "• Nothing pulls the gate back up."),
             size=15, color=DARK)
    add_image(s, FIG_DIR / "gate_reparam.png",
              Inches(6.7), Inches(1.5), w=Inches(6.4))
    add_footer(s, 5, TOTAL)

    # ---------- Slide 6: Fix 1 — gate floor ----------
    s = add_blank_slide(prs)
    add_title(s, "Fix 1: gate floor reparameterization")
    add_image(s, EQ_DIR / "gate_floor.png",
              Inches(1.5), Inches(1.7), w=Inches(10.3))
    add_text(s, Inches(0.6), Inches(3.6), Inches(12.1), Inches(0.6),
             "Bound the effective gate below by a floor. The optimizer can attenuate, never silence.",
             size=18, color=DARK, align=PP_ALIGN.CENTER)
    # three callouts
    cards = [
        ("floor = 0.1", "v2 choice", NAVY),
        ("α_init = 0", "→ g_eff = 0.55 at start", DARK),
        ("0 new params", "purely a reparameterization", GREEN),
    ]
    for i, (a, b, c) in enumerate(cards):
        x = Inches(0.7 + i * 4.15)
        add_text(s, x, Inches(4.7), Inches(4.0), Inches(0.6),
                 a, size=24, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(5.4), Inches(4.0), Inches(0.5),
                 b, size=14, color=LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, 6, TOTAL)

    # ---------- Slide 7: Fix 2 — level mask ----------
    s = add_blank_slide(prs)
    add_title(s, "Fix 2: refine only where small objects live")
    add_image(s, EQ_DIR / "level_mask.png",
              Inches(1.5), Inches(1.7), w=Inches(10.3))
    add_text(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.6),
             "P2 (stride 4) and P3 (stride 8) carry small-object signal. S4/S5 pass through unchanged.",
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    # the four levels visualization
    levels = [
        ("P2", "160×160 = 25,600", "stride 4", True),
        ("P3", "80×80  =  6,400", "stride 8", True),
        ("P4", "40×40  =  1,600", "stride 16", False),
        ("P5", "20×20  =    400", "stride 32", False),
    ]
    for i, (name, count, stride, active) in enumerate(levels):
        x = Inches(0.6 + i * 3.15)
        c = NAVY if active else LIGHT
        bg_color = RGBColor(0xE8, 0xF1, 0xFF) if active else RGBColor(0xF5, 0xF5, 0xF7)
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(4.5), Inches(3.0), Inches(2.0))
        rect.fill.solid(); rect.fill.fore_color.rgb = bg_color
        rect.line.color.rgb = c
        rect.line.width = Pt(2 if active else 0.75)
        rect.shadow.inherit = False
        # clear default text
        rect.text_frame.text = ""
        add_text(s, x, Inches(4.65), Inches(3.0), Inches(0.6),
                 name, size=28, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(5.25), Inches(3.0), Inches(0.4),
                 count, size=12, color=DARK, align=PP_ALIGN.CENTER, font="Menlo")
        add_text(s, x, Inches(5.65), Inches(3.0), Inches(0.4),
                 stride, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(6.05), Inches(3.0), Inches(0.4),
                 "✓ refined" if active else "— skipped",
                 size=12, bold=active, color=c, align=PP_ALIGN.CENTER)
    add_footer(s, 7, TOTAL)

    # ---------- Slide 8: The result ----------
    s = add_blank_slide(prs)
    add_title(s, "Result: the mechanism is now causally measurable")
    add_image(s, EQ_DIR / "v2_pass.png",
              Inches(2.0), Inches(1.4), w=Inches(9.3))
    # the table
    table_data = [
        ["", "ON", "OFF", "Δ (causal)"],
        ["AP",   f"{nums['AP_on']:.2f}",   f"{nums['AP_off']:.2f}",   f"+{nums['AP_on']-nums['AP_off']:.2f}"],
        ["AP_S", f"{nums['APS_on']:.2f}",  f"{nums['APS_off']:.2f}",  f"+{nums['APS_on']-nums['APS_off']:.2f}"],
        ["AP_M", f"{nums['APM_on']:.2f}",  f"{nums['APM_off']:.2f}",  f"+{nums['APM_on']-nums['APM_off']:.2f}"],
        ["AP_L", f"{nums['APL_on']:.2f}",  f"{nums['APL_off']:.2f}",  f"+{nums['APL_on']-nums['APL_off']:.2f}"],
    ]
    rows, cols = len(table_data), len(table_data[0])
    tbl_x, tbl_y = Inches(0.7), Inches(3.4)
    tbl_w, tbl_h = Inches(5.6), Inches(3.0)
    table = s.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, tbl_h).table
    for r, row in enumerate(table_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            tf = cell.text_frame
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Helvetica"
                    run.font.size = Pt(15 if r > 0 else 13)
                    run.font.bold = (r == 0) or (c == 0) or (r == 2 and c == 3)
                    if r == 0:
                        run.font.color.rgb = BG
                    elif r == 2 and c == 3:
                        run.font.color.rgb = GREEN
                    elif c == 3:
                        run.font.color.rgb = GREEN
                    else:
                        run.font.color.rgb = DARK
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif r == 2:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF7, 0xE8)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = BG
    # right side: bar chart
    add_image(s, FIG_DIR / "ablation_bars.png",
              Inches(6.7), Inches(3.4), w=Inches(6.4))
    add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
             f"Same checkpoint, only inference-time toggle. AP_S @ 800 reaches {nums['APS_800']:.2f}.",
             size=12, color=LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, 8, TOTAL)

    # ---------- Slide 9: Trajectory ----------
    s = add_blank_slide(prs)
    add_title(s, "v2 climbs faster and ends higher")
    add_image(s, FIG_DIR / "learning_curves.png",
              Inches(0.7), Inches(1.4), w=Inches(7.3))
    add_image(s, FIG_DIR / "gate_evolution.png",
              Inches(8.2), Inches(1.4), w=Inches(4.9))
    add_text(s, Inches(0.7), Inches(6.3), Inches(7.3), Inches(0.4),
             "v2 first crosses v1's eventual final 33.91 at epoch 8 — four epochs early.",
             size=12, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.2), Inches(6.3), Inches(4.9), Inches(0.4),
             "Gate stays at ≈0.39 (well above the 0.1 floor).",
             size=12, color=LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, 9, TOTAL)

    # ---------- Slide 10: Limitations + repo ----------
    s = add_blank_slide(prs)
    add_title(s, "Limitations and resources")
    # left column: limitations
    add_text(s, Inches(0.6), Inches(1.5), Inches(6.5), Inches(0.5),
             "Limitations",
             size=22, bold=True, color=DARK)
    bullets = [
        ("13 epochs", " instead of the 72 of the standard 6× schedule (compute-limited finetune from public weights)."),
        ("Single seed", " — multi-seed reproducibility study is open."),
        ("Floor & mask not isolated", " — a 2×2 ablation would resolve."),
        ("2× slower than baseline", " — cost is from the P2 stride-4 level, not the feedback module."),
        ("COCO only", " — aerial / surveillance datasets are obvious next steps."),
    ]
    for i, (head, tail) in enumerate(bullets):
        y = Inches(2.15 + i * 0.65)
        # bullet
        b = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.18),
                                Inches(0.10), Inches(0.10))
        b.fill.solid(); b.fill.fore_color.rgb = NAVY
        b.line.fill.background()
        # combined text
        tb = s.shapes.add_textbox(Inches(0.95), y, Inches(6.5), Inches(0.6))
        tf = tb.text_frame
        tf.margin_left = Emu(0); tf.margin_top = Emu(0)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = head
        r1.font.name = "Helvetica"; r1.font.size = Pt(14); r1.font.bold = True
        r1.font.color.rgb = DARK
        r2 = p.add_run(); r2.text = tail
        r2.font.name = "Helvetica"; r2.font.size = Pt(13)
        r2.font.color.rgb = LIGHT
    # right column: hyperlinks (the "interactive" element)
    add_text(s, Inches(7.7), Inches(1.5), Inches(5.4), Inches(0.5),
             "Resources",
             size=22, bold=True, color=DARK)
    add_text(s, Inches(7.7), Inches(2.15), Inches(5.4), Inches(0.4),
             "📦  GitHub repository",
             size=14, bold=True, color=DARK)
    add_hyperlink(s, Inches(7.7), Inches(2.55), Inches(5.4), Inches(0.4),
                  "github.com/HatemSaadallah/feedback-augmented-rtdetr",
                  "https://github.com/HatemSaadallah/feedback-augmented-rtdetr",
                  size=13)
    add_text(s, Inches(7.7), Inches(3.15), Inches(5.4), Inches(0.4),
             "📄  Full 16-page report (PDF)",
             size=14, bold=True, color=DARK)
    add_hyperlink(s, Inches(7.7), Inches(3.55), Inches(5.4), Inches(0.4),
                  "report/main.pdf  ·  view on GitHub",
                  "https://github.com/HatemSaadallah/feedback-augmented-rtdetr/blob/main/report/main.pdf",
                  size=13)
    add_text(s, Inches(7.7), Inches(4.15), Inches(5.4), Inches(0.4),
             "🔬  Reproduce the ablation",
             size=14, bold=True, color=DARK)
    add_hyperlink(s, Inches(7.7), Inches(4.55), Inches(5.4), Inches(0.4),
                  "tools/full_ablation.py  →  --feedback-off",
                  "https://github.com/HatemSaadallah/feedback-augmented-rtdetr/blob/main/rtdetr_pytorch/tools/full_ablation.py",
                  size=13)
    add_text(s, Inches(7.7), Inches(5.15), Inches(5.4), Inches(0.4),
             "🏛  Course",
             size=14, bold=True, color=DARK)
    add_text(s, Inches(7.7), Inches(5.55), Inches(5.4), Inches(0.4),
             "Computer Vision and Image Processing",
             size=13, color=LIGHT)
    add_text(s, Inches(7.7), Inches(5.85), Inches(5.4), Inches(0.4),
             "Bocconi University — MSc AI",
             size=13, color=LIGHT)
    # closer
    add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
             "Thank you. Questions welcome.",
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, 10, TOTAL)

    prs.save(OUT)
    print(f"\nSaved: {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, {TOTAL} slides)")


if __name__ == "__main__":
    generate_equations()
    build()
